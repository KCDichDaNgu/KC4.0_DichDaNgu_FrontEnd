import imp
import io
import os
import re
import aiofiles, asyncio
from docx import Document
import shutil
from sanic.request import File
from docx import Document
from docx.document import Document as _Document
from pptx import Presentation
import openpyxl
from infrastructure.configs.main import GlobalConfig, get_cnf
from docx.table import _Cell, Table
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.text.paragraph import Paragraph
from core.utils.document import check_if_cell_is_string

config: GlobalConfig = get_cnf()
STATIC_FOLDER = config.APP_CONFIG.STATIC_FOLDER

def get_full_path(path: str):

    return f'{STATIC_FOLDER}/{path}'

def extract_file_extension(file_name: str):

    file_name_els = file_name.split('.')

    if len(file_name_els) < 2: return ''

    return file_name.split('.')[-1]

def get_doc_paragraphs(parent):
    if isinstance(parent, _Document):
        parent_elm = parent.element.body
        # print(parent_elm.xml)
    elif isinstance(parent, _Cell):
        parent_elm = parent._tc
    else:
        raise ValueError("something's not right")

    for child in parent_elm.iterchildren():
        if isinstance(child, CT_P):
            yield Paragraph(child, parent)
        elif isinstance(child, CT_Tbl):
            
            table = Table(child, parent)

            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.paragraphs:
                        yield paragraph

def get_doc_file_meta(doc_file: File):

    binary_doc = io.BytesIO(doc_file.body)

    doc = Document(binary_doc)

    doc_paragraphs = list(get_doc_paragraphs(doc))

    total_doc_paragraphs = len(doc_paragraphs)

    full_text = ''

    for paragraph in doc_paragraphs:
        full_text = full_text + paragraph.text

    sentences = re.split('[;.?!]', full_text)

    sentence_count = sum(1 for y in sentences if len(y) > 2)

    return binary_doc, total_doc_paragraphs, sentence_count

def get_presentation_file_meta(presentation_file: File):
    
    binary_presentation = io.BytesIO(presentation_file.body)
    
    presentation = Presentation(binary_presentation)
    
    total_presentation_paragraphs = 0
    total_slides = 0
    slides = [slide for slide in presentation.slides]
    shapes = []
    full_text = ''
    
    for slide in slides:
        total_slides += 1
        for shape in slide.shapes:
            shapes.append(shape)

    for shape in shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame
            for paragraph in text_frame.paragraphs:
                full_text = full_text + paragraph.text
                total_presentation_paragraphs += 1
                
        # if shape.has_table:
        #     for row in shape.table.rows:
        #         for cell in row.cells:
        #             total_presentation_paragraphs += 1
        
    sentences = re.split('[;.?!]', full_text)

    sentence_count = sum(1 for y in sentences if len(y) > 2)
    
    return binary_presentation, total_presentation_paragraphs, total_slides, sentence_count

def get_worksheet_file_meta(worksheet_file: File):
    total_cells = 0
    total_sheets = 0
    full_text = ''
    
    binary_worksheet = io.BytesIO(worksheet_file.body)
    
    worksheet = openpyxl.load_workbook(binary_worksheet)
    
    ws_name_list = worksheet.sheetnames
    
    for ws_name in ws_name_list:
        ws = worksheet[ws_name]
        total_cells += ws.max_column * ws.max_row
        
        for r in range(1,ws.max_row+1):
            for c in range(1,ws.max_column+1):
                cell = ws.cell(r,c)
                
                if check_if_cell_is_string(cell):
                    full_text += str(cell.value) + ';' 
                
    total_sheets = len(ws_name_list)
        
    sentences = re.split('[;.?!]', full_text)

    sentence_count = sum(1 for y in sentences if len(y) > 2)

    return binary_worksheet, total_sheets, total_cells, sentence_count

def get_txt_file_meta(txt_file: File):

    full_text = (txt_file.body.decode('utf-16', errors="ignore"))

    sentences = re.split('[;.?!。]', full_text)

    sentence_count = sum(1 for y in sentences if len(y) > 2)

    return sentence_count

async def delete_files(invalid_file_paths):

    delete_request = []

    for file_path in invalid_file_paths:

        full_file_path = get_full_path(file_path)

        if os.path.exists(full_file_path):
            delete_request.append(aiofiles.os.remove(full_file_path))
    
    await asyncio.gather(*delete_request)

async def delete_folders(invalid_folders):
    for folder in invalid_folders:

        full_file_path = get_full_path(folder)

        if os.path.exists(full_file_path):
            shutil.rmtree(full_file_path)
